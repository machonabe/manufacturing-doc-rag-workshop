# Databricks notebook source
# /// script
# [tool.databricks.environment]
# environment_version = "5"
# ///
# DBTITLE 1,01 セットアップとサンプルデータ生成
# MAGIC %md
# MAGIC # 01_setup_and_sample_data: 環境セットアップと合成データ生成
# MAGIC
# MAGIC ## 目的
# MAGIC - Unity Catalog スキーマ・ Volume を作成
# MAGIC - 製造業の開発ドキュメントを模擬した合成データを生成
# MAGIC
# MAGIC ## 所要時間目安: 10分
# MAGIC
# MAGIC ## 前提
# MAGIC - 00_config が同じフォルダに存在すること
# MAGIC
# MAGIC ## このノートブックで作成されるもの
# MAGIC - スキーマ: `{CATALOG}.{SCHEMA}`
# MAGIC - Volume: `{CATALOG}.{SCHEMA}.{VOLUME}` とそのサブディレクトリ
# MAGIC - 合成ドキュメントファイル 22 件（Excel×6, Word×3, PPT×2, PDF×2, TIFF×3, CSV×6）
# MAGIC
# MAGIC ## ⏱️ 先に実行しておくセル（重要）
# MAGIC Vector Search エンドポイント作成は **10〜20分** かかります。  
# MAGIC スキーマ・ Volume 作成の直後にある「**VS エンドポイント事前作成**」セルを実行したら、完了を待たずに次へ進んでください。
# MAGIC
# MAGIC ## Box からの取り込みを模擬
# MAGIC 実案件では Box MCP サーバー / Box API 経由でファイルを Volume の `raw/` に着地させます。
# MAGIC 本ハンズオンでは、その「取り込み後」の状態をプログラムで生成して再現します。
# MAGIC 詳細は `docs/appendix_box_mcp.md` を参照してください。

# COMMAND ----------

# DBTITLE 1,ライブラリインストール
# Databricks notebook source
# サーバレスコンピュートに必要なライブラリをインストール
# 製造業実務: 非構造化ドキュメントの解析に必要なパッケージ群
%pip install openpyxl python-docx python-pptx reportlab pypdf pillow --quiet

# COMMAND ----------

# DBTITLE 1,Pythonカーネル再起動
# pip install 後はカーネル再起動が必要
# 再起動後、次のセルから自動的に実行が続きます
dbutils.library.restartPython()

# COMMAND ----------

# DBTITLE 1,設定読み込み
# MAGIC %run ./00_config

# COMMAND ----------

# DBTITLE 1,スキーマ・Volume作成
# ==============================================================
# スキーマ・ Volume 作成
# ==============================================================
# 製造業実務: プロジェクト開始時にデータ基盤を整備する工程に相当
spark.sql(f"CREATE SCHEMA IF NOT EXISTS {CATALOG}.{SCHEMA}")
spark.sql(f"CREATE VOLUME IF NOT EXISTS {CATALOG}.{SCHEMA}.{VOLUME}")

import os
for d in [RAW_DIR, IMAGES_DIR, WAVEFORMS_DIR, THUMBS_DIR]:
    os.makedirs(d, exist_ok=True)

print(f"✅ スキーマ作成: {CATALOG}.{SCHEMA}")
print(f"✅ Volume 作成: {VOLUME_PATH}")
print(f"✅ サブディレクトリ: raw/, images/, waveforms/, thumbnails/")

# COMMAND ----------

# DBTITLE 1,⏱️ 時間のかかる処理を先行実行
# MAGIC %md
# MAGIC ## ⏱️ Vector Search エンドポイント事前作成（重要）
# MAGIC
# MAGIC Vector Search エンドポイントの作成には **10〜20分** かかります。  
# MAGIC ここで先に起動しておくことで、02〜04 のノートブックを進めている間にバックグラウンドでプロビジョニングが完了します。
# MAGIC
# MAGIC > **ハンズオン進行のコツ**: このセルを実行したら、完了を待たずに次のセルへ進んでください。  
# MAGIC > エンドポイントの準備完了は `05_vector_index` ノートブックで確認します。

# COMMAND ----------

# DBTITLE 1,VS エンドポイント事前作成（非同期）
# ==============================================================
# Vector Search エンドポイント事前作成（非同期）
# ==============================================================
# 製造業実務: 検索基盤の構築はリードタイムが長いため、
# データ準備と並行してインフラのプロビジョニングを開始する。
from databricks.sdk import WorkspaceClient
from databricks.sdk.service.vectorsearch import EndpointType

w = WorkspaceClient()

# 既存エンドポイントがあればスキップ
try:
    ep = w.vector_search_endpoints.get_endpoint(VS_ENDPOINT_NAME)
    print(f"✅ VS エンドポイント既存: {VS_ENDPOINT_NAME} (状態: {ep.endpoint_status.state.value})")
    print("   → 作成済みのためスキップします")
except Exception:
    print(f"🚀 VS エンドポイント作成開始: {VS_ENDPOINT_NAME}")
    print("   → プロビジョニングに 10〜20分かかります。完了を待たず次へ進んでください。")
    w.vector_search_endpoints.create_endpoint(
        name=VS_ENDPOINT_NAME,
        endpoint_type=EndpointType.STANDARD
    )
    print(f"   → 作成リクエスト送信完了（バックグラウンドで処理中）")
    print(f"   → 状態確認: 05_vector_index ノートブックで確認します")

# COMMAND ----------

# DBTITLE 1,製品マスタデータ定義
# ==============================================================
# 製品マスタデータ
# ==============================================================
# 架空の電子部品メーカー「サンプル電子株式会社」の車載センサー製品
import numpy as np
from io import BytesIO

PRODUCTS = [
    {"id": "SNS-100", "name": "高精度温度センサー", "spec_id": "SPEC-SNS-100-v2",
     "temp_range": "-40~125\u00b0C", "accuracy": "\u00b10.5\u00b0C", "response": "200ms"},
    {"id": "SNS-200", "name": "車載加速度センサー", "spec_id": "SPEC-SNS-200-v1",
     "temp_range": "-40~150\u00b0C", "accuracy": "\u00b12%FS", "response": "50ms"},
    {"id": "SNS-300", "name": "圧力センサーモジュール", "spec_id": "SPEC-SNS-300-v1",
     "temp_range": "-30~120\u00b0C", "accuracy": "\u00b11%FS", "response": "100ms"},
]

# COMMAND ----------

# DBTITLE 1,波形生成ユーティリティ
# ==============================================================
# 波形データ生成ユーティリティ
# ==============================================================
# 製造業実務: センサーのステップ応答・振動・パルス応答などの試験波形を模擬
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

def gen_wave(wtype="step", noise=0.05, n=200):
    """試験波形データを生成"""
    t = np.linspace(0, 1, n)
    if wtype == "step":
        s = np.where(t > 0.2, 1.0 - np.exp(-(t - 0.2) * 10), 0.0)
    elif wtype == "sine":
        s = np.sin(2 * np.pi * 5 * t) * np.exp(-t * 2)
    elif wtype == "pulse":
        s = np.exp(-((t - 0.5) ** 2) / 0.01)
    else:
        s = np.random.randn(n) * 0.1
    return t, s + np.random.randn(n) * noise

def chart_png(t, s, title):
    """波形チャートをPNGバイト列として生成"""
    fig, ax = plt.subplots(figsize=(4, 2.5), dpi=100)
    ax.plot(t * 1000, s, 'b-', lw=0.8)
    ax.set_xlabel("Time [ms]"); ax.set_ylabel("Voltage [V]")
    ax.set_title(title, fontsize=9); ax.grid(True, alpha=0.3)
    buf = BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight')
    plt.close(fig); buf.seek(0)
    return buf

def save_wb(wb, path):
    """openpyxl Workbook を Volume に保存（FUSE互換）"""
    buf = BytesIO(); wb.save(buf)
    with open(path, 'wb') as f:
        f.write(buf.getvalue())

print("✅ ユーティリティ関数定義完了")

# COMMAND ----------

# DBTITLE 1,Excelファイル生成（3レイアウト）
# ==============================================================
# Excel 試験成績書生成（6ファイル、レイアウトA/B/Cの3種）
# ==============================================================
# 製造業実務: 試験成績書は部署・時期によってフォーマットが異なる。
# 意図的に3種のレイアウトで生成し、後段の抽出処理で「不統一フォーマットへの対応」を学ぶ。
from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.drawing.image import Image as XlImage
import os, csv

def make_layout_a(p, tid, path):
    """レイアウトA: 標準型 3シート（試験条件/測定結果/波形）"""
    wb = Workbook()
    # シート1: 試験条件
    ws = wb.active; ws.title = "試験条件"
    ws['A1'] = f"試験成績書 - {p['name']}"; ws['A1'].font = Font(bold=True, size=14)
    ws.merge_cells('A1:E1')
    ws['A2'] = f"製品ID: {p['id']}"; ws['A3'] = f"仕様書: {p['spec_id']}"
    ws['A4'] = f"試験ID: {tid}"; ws['A5'] = "試験日: 2026-07-15"
    for c, h in enumerate(["項目","条件","規格値","単位","備考"], 1):
        ws.cell(7, c, h).font = Font(bold=True)
    for r, d in enumerate([["動作温度","25\u00b0C",p['temp_range'],"\u00b0C","常温"],
                           ["供給電圧","5.0V","4.5~5.5V","V","定格"],
                           ["負荷抵抗","10k\u03a9",">1k\u03a9","\u03a9",""]], 8):
        for c, v in enumerate(d, 1): ws.cell(r, c, v)
    # シート2: 測定結果（取り消し線あり）
    ws2 = wb.create_sheet("測定結果")
    ws2['A1'] = f"測定結果 - {tid}"; ws2['A2'] = f"製品: {p['id']}"
    for c, h in enumerate(["測定項目","旧測定値","修正値","判定"], 1):
        ws2.cell(4, c, h).font = Font(bold=True)
    ws2.cell(5,1,"応答時間"); ws2.cell(5,2,"210ms"); ws2['B5'].font=Font(strike=True)
    ws2.cell(5,3,"195ms"); ws2.cell(5,4,"PASS")
    ws2.cell(6,1,"オフセット電圧"); ws2.cell(6,2,"0.15mV"); ws2['B6'].font=Font(strike=True)
    ws2.cell(6,3,"0.08mV"); ws2.cell(6,4,"PASS")
    # シート3: 波形データ + 埋め込みチャート
    ws3 = wb.create_sheet("波形データ")
    ws3['A1'] = f"波形 - {tid} - {p['id']}"
    t, s = gen_wave("step", 0.03)
    ws3['A3'] = "time_ms"; ws3['B3'] = "voltage_V"
    for i, (tv, sv) in enumerate(zip(t, s), 4):
        ws3.cell(i, 1, round(tv*1000, 3)); ws3.cell(i, 2, round(sv, 6))
    ws3.add_image(XlImage(chart_png(t, s, f"Step Response - {p['id']}")), 'D3')
    save_wb(wb, path)

def make_layout_b(p, tid, path):
    """レイアウトB: 結合セル多用・総合試験形式"""
    wb = Workbook()
    ws = wb.active; ws.title = "総合試験"
    ws.merge_cells('A1:F1'); ws['A1']=f"【{p['id']}】総合試験報告書"
    ws['A1'].font = Font(bold=True, size=14)
    ws.merge_cells('A2:C2'); ws['A2'] = f"仕様書: {p['spec_id']}"
    ws.merge_cells('D2:F2'); ws['D2'] = f"試験番号: {tid}"
    for c, h in enumerate(["No.","試験項目","条件","測定値","合否","備考"], 1):
        ws.cell(4, c, h).font = Font(bold=True)
    for r, d in enumerate([["1","耐熱","85\u00b0C/1000h","異常なし","PASS",""],
                           ["2","耐湿","85\u00b0C/85%RH","異常なし","PASS",""],
                           ["3","熱衝撃","\u221240~125\u00b0C","異常なし","PASS",""]], 5):
        for c, v in enumerate(d, 1): ws.cell(r, c, v)
    ws.cell(8,1,"4"); ws.cell(8,2,"応答速度"); ws.cell(8,3,"常温")
    ws.cell(8,4,"250ms"); ws['D8'].font = Font(strike=True)
    ws.cell(8,5,"再測定"); ws.cell(8,6,"\u2192190ms PASS")
    ws2 = wb.create_sheet("波形記録")
    t, s = gen_wave("sine", 0.02)
    ws2['A1'] = f"振動波形 - {p['id']} - {tid}"
    ws2['A2'] = "time_ms"; ws2['B2'] = "accel_g"
    for i, (tv, sv) in enumerate(zip(t, s), 3):
        ws2.cell(i, 1, round(tv*1000, 3)); ws2.cell(i, 2, round(sv, 6))
    ws2.add_image(XlImage(chart_png(t, s, f"Vibration - {p['id']}")), 'D2')
    save_wb(wb, path)

def make_layout_c(products, tid, path):
    """レイアウトC: 複数製品を 1 ファイルに集約（製品ID別シート）"""
    wb = Workbook()
    for p in products:
        ws = wb.create_sheet(p['id'])
        ws['A1'] = f"環境試験 - {p['id']} - {p['name']}"
        ws['A1'].font = Font(bold=True)
        ws['A2'] = f"SPEC: {p['spec_id']}"; ws['A3'] = f"TEST: {tid}"
        ws['A5'] = "温度[\u00b0C]"; ws['B5'] = "出力[V]"; ws['C5'] = "判定"
        for i, temp in enumerate([-40, -20, 0, 25, 50, 85, 125], 6):
            ws.cell(i, 1, temp)
            ws.cell(i, 2, round(2.5 + np.random.randn() * 0.01, 4))
            ws.cell(i, 3, "PASS")
    if "Sheet" in wb.sheetnames: del wb["Sheet"]
    save_wb(wb, path)

# --- Excel ファイル生成実行 ---
excel_files = []; tid_n = 1
for p in PRODUCTS[:2]:
    tid = f"TEST-2026-{tid_n:03d}"; fname = f"TEST_{p['id']}_{tid}_layoutA.xlsx"
    make_layout_a(p, tid, os.path.join(RAW_DIR, fname))
    excel_files.append(fname); tid_n += 1
for p in PRODUCTS[1:]:
    tid = f"TEST-2026-{tid_n:03d}"; fname = f"TEST_{p['id']}_{tid}_layoutB.xlsx"
    make_layout_b(p, tid, os.path.join(RAW_DIR, fname))
    excel_files.append(fname); tid_n += 1
tid = f"TEST-2026-{tid_n:03d}"; fname = f"TEST_MULTI_{tid}_layoutC.xlsx"
make_layout_c(PRODUCTS, tid, os.path.join(RAW_DIR, fname))
excel_files.append(fname); tid_n += 1
tid = f"TEST-2026-{tid_n:03d}"; fname = f"TEST_{PRODUCTS[2]['id']}_{tid}_layoutA.xlsx"
make_layout_a(PRODUCTS[2], tid, os.path.join(RAW_DIR, fname))
excel_files.append(fname); tid_n += 1

print(f"✅ Excel 試験成績書 {len(excel_files)} ファイル生成完了")
for f in excel_files: print(f"   {f}")

# COMMAND ----------

# DBTITLE 1,Word/PPT/PDF/TIFF/CSV生成
# ==============================================================
# Word 仕様書 / PowerPoint / PDF / TIFF / CSV 生成
# ==============================================================
from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas as pdf_canvas
from PIL import Image, ImageDraw

# --- Word 仕様書 ×3 ---
def create_word_spec(p, path):
    """製造業実務: 製品仕様書は設計・品質・試験の根拠ドキュメント"""
    doc = Document()
    doc.add_heading(f"製品仕様書 - {p['name']}", 0)
    doc.add_paragraph(f"仕様書番号: {p['spec_id']}")
    doc.add_paragraph(f"製品ID: {p['id']}")
    doc.add_paragraph(f"発行日: 2026-04-01  改訂: Rev.2")
    doc.add_heading("1. 概要", level=1)
    doc.add_paragraph(
        f"{p['name']}は、車載環境向けに設計された高信頼性センサーです。"
        f"動作温度範囲 {p['temp_range']} での安定動作を保証します。")
    doc.add_heading("2. 電気的特性", level=1)
    table = doc.add_table(rows=5, cols=4); table.style = 'Table Grid'
    for i, h in enumerate(["項目","最小","標準","最大"]):
        table.rows[0].cells[i].text = h
    for r, row in enumerate([["\u4f9b\u7d66\u96fb\u5727","4.5V","5.0V","5.5V"],
                             ["消費電流","\u2014","5mA","10mA"],
                             ["精度","\u2014",p['accuracy'],"\u2014"],
                             ["応答時間","\u2014",p['response'],"500ms"]], 1):
        for c, v in enumerate(row): table.rows[r].cells[c].text = v
    doc.add_heading("3. 動作温度範囲", level=1)
    doc.add_paragraph(f"動作温度: {p['temp_range']}")
    doc.add_heading("4. 関連ドキュメント", level=1)
    doc.add_paragraph(f"・ブロック図: BLK_{p['id']}_rev2.tiff")
    doc.add_paragraph(f"・試験成績書: TEST-2026-xxx")
    buf = BytesIO(); doc.save(buf)
    with open(path, 'wb') as f: f.write(buf.getvalue())

word_files = []
for p in PRODUCTS:
    fname = f"SPEC_{p['id']}_v2.docx"
    create_word_spec(p, os.path.join(RAW_DIR, fname))
    word_files.append(fname)
print(f"✅ Word {len(word_files)} ファイル")

# --- PowerPoint ×2 ---
def create_pptx(title, products_subset, path):
    """製造業実務: 設計レビュー資料"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "サンプル電子株式会社 設計部\n2026年7月"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "レビュー対象"
    body = slide2.placeholders[1]
    body.text = "\n".join([f"・{p['id']}: {p['name']} (SPEC: {p['spec_id']})" for p in products_subset])
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "課題と対応"
    slide3.placeholders[1].text = (
        "1. SNS-200 応答時間改善\n"
        "2. SNS-300 耐圧マージン確保\n"
        "3. EMC試験スケジュール調整")
    buf = BytesIO(); prs.save(buf)
    with open(path, 'wb') as f: f.write(buf.getvalue())

pptx_files = []
create_pptx("2026Q3 設計レビュー", PRODUCTS, os.path.join(RAW_DIR, "REVIEW_2026Q3_design.pptx"))
pptx_files.append("REVIEW_2026Q3_design.pptx")
create_pptx("SNS-200 改善検討", [PRODUCTS[1]], os.path.join(RAW_DIR, "REVIEW_SNS200_improvement.pptx"))
pptx_files.append("REVIEW_SNS200_improvement.pptx")
print(f"✅ PPT {len(pptx_files)} ファイル")

# --- PDF ×2 ---
def create_pdf(p, path):
    """製造業実務: データシート（英文PDF）"""
    c = pdf_canvas.Canvas(path, pagesize=A4); w, h = A4
    c.setFont("Helvetica-Bold", 18)
    c.drawString(50, h-50, f"Datasheet: {p['id']}")
    c.setFont("Helvetica", 12)
    c.drawString(50, h-80, f"Product: {p['name']}")
    c.drawString(50, h-100, f"Spec: {p['spec_id']}")
    c.drawString(50, h-130, "Electrical Characteristics")
    c.setFont("Helvetica", 10)
    y = h - 160
    for line in [f"Supply Voltage: 4.5 ~ 5.5V",
                 f"Operating Temp: {p['temp_range']}",
                 f"Accuracy: {p['accuracy']}",
                 f"Response Time: {p['response']}",
                 f"Package: QFN-16 (4x4mm)",
                 f"Related: BLK_{p['id']}_rev2.tiff"]:
        c.drawString(70, y, line); y -= 20
    c.save()

pdf_files = []
for p in PRODUCTS[:2]:
    fname = f"DS_{p['id']}_datasheet.pdf"
    create_pdf(p, os.path.join(RAW_DIR, fname))
    pdf_files.append(fname)
print(f"✅ PDF {len(pdf_files)} ファイル")

# --- TIFF ブロック図 ×3 ---
def create_block_diagram(p, path):
    """製造業実務: ブロック図は製品構造の理解に不可欠"""
    img = Image.new('RGB', (800, 400), 'white')
    draw = ImageDraw.Draw(img)
    draw.text((20, 10), f"Block Diagram - {p['id']} ({p['name']})", fill='black')
    blocks = [("Sensor\nElement",50,100),("AFE",250,100),("ADC",450,100),("DSP",650,100)]
    for label, x, y in blocks:
        draw.rectangle([x, y, x+140, y+80], outline='black', width=2)
        draw.text((x+20, y+25), label, fill='black')
    for i in range(len(blocks)-1):
        x1 = blocks[i][1]+140; x2 = blocks[i+1][1]
        draw.line([(x1, 140), (x2, 140)], fill='black', width=2)
    draw.text((20, 250), f"SPEC: {p['spec_id']}", fill='gray')
    draw.text((20, 270), f"Output: I2C/SPI | Power: 5V | Temp: {p['temp_range']}", fill='gray')
    img.save(path, format='TIFF')

tiff_files = []
for p in PRODUCTS:
    fname = f"BLK_{p['id']}_rev2.tiff"
    create_block_diagram(p, os.path.join(RAW_DIR, fname))
    tiff_files.append(fname)
print(f"✅ TIFF {len(tiff_files)} ファイル")

# --- 波形 CSV ×6 ---
csv_files = []; wave_types = ["step","sine","pulse","step","sine","pulse"]
for i, (p, wt) in enumerate(zip(PRODUCTS * 2, wave_types)):
    tid = f"TEST-2026-{i+1:03d}"
    fname = f"WAVE_{p['id']}_{tid}.csv"
    t, s = gen_wave(wt, 0.02)
    with open(os.path.join(RAW_DIR, fname), 'w', newline='') as f:
        writer = csv.writer(f)
        writer.writerow(["time_ms", "voltage_V"])
        for tv, sv in zip(t, s):
            writer.writerow([round(tv*1000, 3), round(sv, 6)])
    csv_files.append(fname)
print(f"✅ CSV {len(csv_files)} ファイル")

# COMMAND ----------

# DBTITLE 1,検証: 生成ファイル一覧と件数確認
# ==============================================================
# 検証: 生成ファイルの確認
# ==============================================================
all_files = sorted(os.listdir(RAW_DIR))
print(f"\n{'='*60}")
print(f"Volume raw/ の内容 ({len(all_files)} ファイル):")
print(f"{'='*60}")
for f in all_files:
    size = os.path.getsize(os.path.join(RAW_DIR, f))
    print(f"  {f:50s} {size:>8,} bytes")

# アサーション
assert len([f for f in all_files if f.endswith('.xlsx')]) >= 6, "Excelファイル数不足"
assert len([f for f in all_files if f.endswith('.docx')]) >= 3, "Wordファイル数不足"
assert len([f for f in all_files if f.endswith('.tiff')]) >= 3, "TIFFファイル数不足"
assert len([f for f in all_files if f.endswith('.csv')]) >= 6, "CSVファイル数不足"
assert len(all_files) >= 20, f"合計ファイル数不足: {len(all_files)}"

print(f"\n✅ 01_setup_and_sample_data 完了: {len(all_files)} ファイル生成済み")
