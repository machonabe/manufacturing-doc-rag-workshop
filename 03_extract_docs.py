# Databricks notebook source
# DBTITLE 1,03 ドキュメント抽出: 紹介
# MAGIC %md
# MAGIC # 03_extract_docs: Word/PPT/PDF/TIFF/CSV からのテキスト抽出
# MAGIC
# MAGIC ## 目的
# MAGIC - Word 仕様書からテキスト・表を抽出
# MAGIC - PowerPoint レビュー資料からスライドテキストを抽出
# MAGIC - PDF データシートからテキストを抽出
# MAGIC - TIFF ブロック図をサムネイル化し media_assets に登録
# MAGIC - 波形 CSV を media_assets に登録
# MAGIC
# MAGIC ## 所要時間目安: 10分
# MAGIC
# MAGIC ## 前提
# MAGIC - 01, 02 ノートブックが実行済み
# MAGIC
# MAGIC ## このノートブックで更新されるもの
# MAGIC - `documents` テーブルに Word/PPT/PDF/TIFF/CSV 分を追加
# MAGIC - `media_assets` テーブルにブロック図・波形を追加
# MAGIC - `thumbnails/` に TIFF のサムネイル PNG

# COMMAND ----------

# DBTITLE 1,設定読み込み
# MAGIC %run ./00_config

# COMMAND ----------

# DBTITLE 1,ドキュメント抽出メインロジック
# ==============================================================
# Word / PPT / PDF / TIFF / CSV からのテキスト抽出
# ==============================================================
import os, re, uuid
from datetime import datetime
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from PIL import Image
from pyspark.sql import Row

def extract_ids(text):
    product_ids = list(set(re.findall(r'SNS-\d{3}', str(text))))
    spec_ids = list(set(re.findall(r'SPEC-SNS-\d{3}-v\d+', str(text))))
    return product_ids, spec_ids

doc_records = []
media_records = []

# --- Word 仕様書 ---
print("=== Word 仕様書抽出 ===")
for fname in sorted(f for f in os.listdir(RAW_DIR) if f.endswith('.docx')):
    fpath = os.path.join(RAW_DIR, fname)
    doc = Document(fpath)
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            text_parts.append(row_text)
    full_text = "\n".join(text_parts)
    prod_ids, spec_ids = extract_ids(full_text)
    doc_records.append(Row(
        doc_id=fname, file_name=fname, file_path=fpath,
        doc_type="word", title=f"Word仕様書 - {fname}",
        product_ids=prod_ids, spec_ids=spec_ids,
        keywords=[], summary=full_text[:300],
        ingested_at=datetime.now().isoformat()
    ))
    print(f"  ✔ {fname} ({len(full_text)}文字, 製品: {prod_ids})")

# --- PowerPoint ---
print("\n=== PowerPoint 抽出 ===")
for fname in sorted(f for f in os.listdir(RAW_DIR) if f.endswith('.pptx')):
    fpath = os.path.join(RAW_DIR, fname)
    prs = Presentation(fpath)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
    full_text = "\n".join(text_parts)
    prod_ids, spec_ids = extract_ids(full_text)
    doc_records.append(Row(
        doc_id=fname, file_name=fname, file_path=fpath,
        doc_type="pptx", title=f"PPTレビュー資料 - {fname}",
        product_ids=prod_ids, spec_ids=spec_ids,
        keywords=[], summary=full_text[:300],
        ingested_at=datetime.now().isoformat()
    ))
    print(f"  ✔ {fname} ({len(full_text)}文字, 製品: {prod_ids})")

# --- PDF ---
print("\n=== PDF 抽出 ===")
for fname in sorted(f for f in os.listdir(RAW_DIR) if f.endswith('.pdf')):
    fpath = os.path.join(RAW_DIR, fname)
    reader = PdfReader(fpath)
    text_parts = [page.extract_text() or "" for page in reader.pages]
    full_text = "\n".join(text_parts)
    prod_ids, spec_ids = extract_ids(full_text)
    doc_records.append(Row(
        doc_id=fname, file_name=fname, file_path=fpath,
        doc_type="pdf", title=f"PDFデータシート - {fname}",
        product_ids=prod_ids, spec_ids=spec_ids,
        keywords=[], summary=full_text[:300],
        ingested_at=datetime.now().isoformat()
    ))
    print(f"  ✔ {fname} ({len(full_text)}文字, 製品: {prod_ids})")

# --- TIFF ブロック図 ---
print("\n=== TIFF ブロック図 ===")
for fname in sorted(f for f in os.listdir(RAW_DIR) if f.endswith('.tiff')):
    fpath = os.path.join(RAW_DIR, fname)
    # サムネイル生成
    img = Image.open(fpath)
    thumb_fname = fname.replace('.tiff', '_thumb.png')
    thumb_path = os.path.join(THUMBS_DIR, thumb_fname)
    img.thumbnail((400, 200))
    img.save(thumb_path, 'PNG')
    
    prod_ids, _ = extract_ids(fname)
    asset_id = str(uuid.uuid4())[:8]
    media_records.append(Row(
        asset_id=asset_id, asset_type="block_diagram",
        file_path=fpath, thumbnail_path=thumb_path,
        source_doc_id=fname, source_file_name=fname,
        source_sheet="", anchor_cell="",
        product_id=prod_ids[0] if prod_ids else "",
        spec_id="", test_id="",
        description=f"ブロック図 - {fname}"
    ))
    doc_records.append(Row(
        doc_id=fname, file_name=fname, file_path=fpath,
        doc_type="tiff", title=f"ブロック図 - {fname}",
        product_ids=prod_ids, spec_ids=[],
        keywords=[], summary=f"製品構造ブロック図 ({', '.join(prod_ids)})",
        ingested_at=datetime.now().isoformat()
    ))
    print(f"  ✔ {fname} → サムネイル: {thumb_fname}")

# --- 波形 CSV ---
print("\n=== 波形 CSV ===")
for fname in sorted(f for f in os.listdir(RAW_DIR) if f.startswith('WAVE_') and f.endswith('.csv')):
    fpath = os.path.join(RAW_DIR, fname)
    prod_ids, _ = extract_ids(fname)
    test_ids = re.findall(r'TEST-\d{4}-\d{3}', fname)
    asset_id = str(uuid.uuid4())[:8]
    media_records.append(Row(
        asset_id=asset_id, asset_type="waveform_csv",
        file_path=fpath, thumbnail_path="",
        source_doc_id=fname, source_file_name=fname,
        source_sheet="", anchor_cell="",
        product_id=prod_ids[0] if prod_ids else "",
        spec_id="", test_id=test_ids[0] if test_ids else "",
        description=f"波形データ - {fname}"
    ))
    doc_records.append(Row(
        doc_id=fname, file_name=fname, file_path=fpath,
        doc_type="csv", title=f"波形CSV - {fname}",
        product_ids=prod_ids, spec_ids=[],
        keywords=[], summary=f"波形生データ ({', '.join(prod_ids)})",
        ingested_at=datetime.now().isoformat()
    ))
    print(f"  ✔ {fname}")

print(f"\n抽出完了: documents {len(doc_records)}件, media_assets {len(media_records)}件")

# COMMAND ----------

# DBTITLE 1,Deltaテーブルへの追記
# ==============================================================
# documents テーブルへの追記（Excel分に非 Excel ドキュメントを追加）
# ==============================================================
if doc_records:
    df_new_docs = spark.createDataFrame(doc_records)
    df_new_docs.write.mode("append").saveAsTable(FQ_DOCUMENTS)
    print(f"✅ {FQ_DOCUMENTS}: {len(doc_records)} 件追加")

if media_records:
    df_new_media = spark.createDataFrame(media_records)
    df_new_media.write.mode("append").saveAsTable(FQ_MEDIA_ASSETS)
    print(f"✅ {FQ_MEDIA_ASSETS}: {len(media_records)} 件追加")

# 統計確認
total_docs = spark.sql(f"SELECT COUNT(*) as cnt FROM {FQ_DOCUMENTS}").first().cnt
total_media = spark.sql(f"SELECT COUNT(*) as cnt FROM {FQ_MEDIA_ASSETS}").first().cnt
print(f"\n統計: documents 合計 {total_docs}件, media_assets 合計 {total_media}件")

# COMMAND ----------

# DBTITLE 1,検証
# ==============================================================
# 検証
# ==============================================================
# ドキュメント種別の内訳確認
print("ドキュメント種別内訳:")
spark.sql(f"SELECT doc_type, COUNT(*) as cnt FROM {FQ_DOCUMENTS} GROUP BY doc_type ORDER BY doc_type").show()

# メディア種別の内訳確認
print("メディアアセット種別内訳:")
spark.sql(f"SELECT asset_type, COUNT(*) as cnt FROM {FQ_MEDIA_ASSETS} GROUP BY asset_type ORDER BY asset_type").show()

# サムネイル確認
thumbs = [f for f in os.listdir(THUMBS_DIR) if f.endswith('.png')]
assert len(thumbs) >= 3, f"サムネイル数不足: {len(thumbs)}"
print(f"サムネイル: {len(thumbs)} 件")

print("\n✅ 03_extract_docs 完了")
